# utils.py
import os
from github import Github, GithubException
import difflib
import re
from docx import Document as DocxDocument
import pdfplumber
from pptx import Presentation
import io
import requests
from requests.auth import HTTPBasicAuth
import base64

sampling_params = {
    "temperature": 0.2,       # More deterministic responses for precision
    "top_p": 0.9,             # Allows for diverse but coherent outputs
    "n": 1,                   # Generate one review
    "frequency_penalty": 0.7, # Stronger reduction of repetition
    "presence_penalty": 0.4,  # Encourages introduction of new ideas
}

def detect_vulnerabilities(code):
    vulnerabilities = []

    # Example checks for vulnerabilities
    if "eval(" in code or "exec(" in code:
        vulnerabilities.append("Avoid using eval() or exec() as they can lead to code injection vulnerabilities.")

    if re.search(r"(\bSELECT\b|\bUPDATE\b|\bDELETE\b|\bINSERT\b)\s+\w+\s+FROM\s+\w+\s*;", code, re.IGNORECASE):
        vulnerabilities.append("Possible SQL Injection vulnerability detected. Use parameterized queries instead.")

    # Check for common XSS patterns
    if re.search(r"<script.*?>", code, re.IGNORECASE):
        vulnerabilities.append("Potential Cross-Site Scripting (XSS) vulnerability detected. Avoid using inline scripts.")

    # Check for hardcoded credentials (this is a simplistic check)
    if re.search(r"(password|secret|api_key)\s*=\s*['\"].*['\"]", code, re.IGNORECASE):
        vulnerabilities.append("Hardcoded credentials detected. Consider using environment variables or a secure vault.")

    return vulnerabilities

def determine_severity_from_score(total_score):
    if total_score == 0:
        color = "darkgreen"
        message = "No errors, perfect!"
    elif total_score <= 10:
        color = "lightgreen"
        message = "Low severity"
    elif 11 <= total_score <= 20:
        color = "yellow"
        message = "Medium severity"
    elif 21 <= 30:
        color = "orange"
        message = "High severity"
    else:
        color = "red"
        message = "Critical severity"

    return {
        "color": color,
        "severity_message": message
    }


error_weights = {
    "Syntax Errors": 2,
    "Runtime Errors": 2,
    "Logical Errors": 1,
    "Validation Errors": 1,
    "Compilation Errors": 2,
}


def severity(error_count, code_to_review):
    # Calculate total severity score
    total_score = calculate_severity(error_count)

    # Determine severity level based on total score
    severity_result = determine_severity_from_score(total_score)
    color = severity_result['color']
    severity_message = severity_result['severity_message']

    # Detect vulnerabilities in the code
    vulnerabilities_found = detect_vulnerabilities(code_to_review)

    return {
        "total_score": total_score,
        "color": color,
        "severity_message": severity_message,
        "total_errors": sum(error_count.values()),  # Total number of errors
        "vulnerabilities": vulnerabilities_found,  # List of detected vulnerabilities
    }


def calculate_severity(error_count):
    total_score = 0

    # Iterate over the error_count dictionary
    for error_type, count in error_count.items():
        # Check if the error_type exists in the predefined weights
        if error_type in error_weights:
            # Multiply the count by the weight of the error type and add it to the total score
            total_score += count * error_weights[error_type]

    return total_score
import pandas as pd

def generate_diff_dataframe(old_content, new_content):
    # Use difflib to generate the unified diff
    diff_lines = list(difflib.unified_diff(old_content, new_content, lineterm=''))

    # Track line numbers in old and new content
    old_line_num = 1
    new_line_num = 1
    diff_data = []

    for line in diff_lines:
        if line.startswith('---') or line.startswith('+++'):
            continue  # Skip version header lines
        elif line.startswith('@@'):
            parts = line.split(' ')
            old_line_num = int(parts[1].split(',')[0].replace('-', ''))
            new_line_num = int(parts[2].split(',')[0].replace('+', ''))
        elif line.startswith('-'):
            # Line removed from old content
            diff_data.append({
                'Old Line': f'{old_line_num}: {line[1:]}',  # Strip the '-'
                'New Line': '',
                'Change': 'Removed'
            })
            old_line_num += 1
        elif line.startswith('+'):
            # Line added to new content
            diff_data.append({
                'Old Line': '',
                'New Line': f'{new_line_num}: {line[1:]}',  # Strip the '+'
                'Change': 'Added'
            })
            new_line_num += 1
        elif line.startswith(' '):
            # Unchanged line
            old_line_num += 1
            new_line_num += 1

    # Create a DataFrame from the diff data
    df_diff = pd.DataFrame(diff_data)

    # Scan for matching line numbers in Old Line and New Line columns
    for i, row in df_diff.iterrows():
        old_line = row['Old Line'].split(':')[0] if row['Old Line'] else None
        for j, inner_row in df_diff.iterrows():
            new_line = inner_row['New Line'].split(':')[0] if inner_row['New Line'] else None
            # If old line number matches new line number, combine them
            if old_line and new_line and old_line == new_line:
                # Combine the two rows
                df_diff.at[i, 'New Line'] = inner_row['New Line']
                df_diff.at[i, 'Change'] = 'Modified'
                # Drop the row that was combined
                df_diff.drop(j, inplace=True)
                break  # Exit inner loop since we found a match

    # Reset the index after dropping rows
    df_diff.reset_index(drop=True, inplace=True)

    # Convert DataFrame to JSON for the frontend (without styling)
    diff_json = df_diff.to_dict(orient='records')

    return diff_json


def calculate_score(org_std_text, new_code, client, model_type):
    # Define prompts based on organizational code standards
    explain_prompt = (
        "Evaluate the following code according to the organization standards. "
        "Provide an overall score out of 10 in the format 'Overall Score: (score)/10' and concisely explain the reasoning behind your score.\n\n"
        f"Code:\n{new_code}\n\n"
        f"Organization Standards:\n{org_std_text}"
    )
    explain_response = client.chat.completions.create(
        messages=[{"role": "user", "content": explain_prompt}],
        model=model_type, **sampling_params
    )

    explanation = explain_response.choices[0].message.content

    # Assuming the response is a score as a float, parse it
    score_match = re.search(r"Overall Score.*?(\d+(\.\d+)?)/10",
                            explanation)  # This will match a number (integer or float)

    if score_match:
        score = float(score_match.group(1))  # Extract the score and convert it to a float
    else:
        return {
            "score": None,
            "explanation": "Failed to parse score",
            "color": "red",
            "message": "Failed to parse score"
        }

    # Determine the score message based on the score
    if score >= 9:
        color = "green"
        message = "Excellent"
    elif score >= 7:
        color = "yellow"
        message = "Good"
    elif score >= 4:
        color = "orange"
        message = "Average"
    else:
        color = "red"
        message = "Poor"

    # Return score and message for React to display
    return {
        "score": score,
        "explanation": explanation,
        "color": color,
        "message": message,
    }

import git


def process_folder_or_repo(input_path):
    # Check if input_path is a GitHub repository URL
    if "github.com" in input_path:
        # Clone the GitHub repository
        repo_name = input_path.split("/")[-1].replace(".git", "")
        if not os.path.exists(repo_name):
            try:
                git.Repo.clone_from(input_path, repo_name)
            except Exception as e:
                return [], f"Error cloning the repository: {str(e)}"
        folder_path = repo_name
    else:
        # Input is a local folder path
        folder_path = input_path

    # Check if the folder exists
    if os.path.exists(folder_path):
        code_files = []  # Initialize list to store valid code files

        # Use os.walk to recursively search through subfolders and files
        for root, dirs, files in os.walk(folder_path):
            for file in files:
                # Check if the file has one of the allowed extensions
                if file.endswith(('.py', '.js', '.java', '.html', '.css', '.cpp')):
                    code_files.append(os.path.join(root, file))  # Add full path of the file

        if code_files:
            return code_files, folder_path
        else:
            return [], "No code files found in the provided folder or repository."
    else:
        return [], "Provided folder path does not exist."


import traceback
from github import Github, GithubException


def process_pull_request(token, repo_name, pull_number):
    """
    Fetches the pull request object from the GitHub API.
    """
    try:
        # Initialize GitHub client
        g = Github(token)
        user = g.get_user()
        owner = user.login  # Owner inferred from the token

        # Fetch repository
        try:
            repo = g.get_repo(f"{owner}/{repo_name}")
        except GithubException as ge:
            print(f"GithubException (repo fetch): {traceback.format_exc()}")
            raise ValueError(
                f"Failed to fetch repository {repo_name} for user {owner}. Ensure the repository exists and is accessible.")

        # Fetch pull request
        try:
            pull_request = repo.get_pull(int(pull_number))
        except GithubException as ge:
            print(f"GithubException (PR fetch): {traceback.format_exc()}")
            raise ValueError(f"Failed to fetch PR {pull_number} from {repo_name}. Make sure the PR number is correct.")

        return pull_request

    except ValueError as ve:
        print(f"ValueError: {traceback.format_exc()}")
        raise ValueError(f"Invalid input: {ve}")
    except GithubException as ge:
        print(f"GithubException: {traceback.format_exc()}")
        raise ValueError(f"GitHub API error: {ge.data.get('message', 'Unknown error')} (HTTP {ge.status})")
    except Exception as e:
        print(f"Unexpected error: {traceback.format_exc()}")
        raise ValueError(f"Unexpected error: {str(e)}")


def extract_changed_code(old_code, new_code):
    """Extract modified lines and clean them for review."""
    old_code_lines = old_code.splitlines(keepends=True)
    new_code_lines = new_code.splitlines(keepends=True)

    # Calculate the diff
    diff = list(difflib.unified_diff(old_code_lines, new_code_lines, lineterm=''))
    return diff  # Return the raw diff output
import re
from groq import Groq  # Assuming you're using OpenAI API for error detection and suggestions


def load_documents_from_files(files):
    documents = {}

    # If `files` is a single file, make it iterable
    if not isinstance(files, list):
        files = [files]

    for file in files:
        # If file is passed as bytes without the .name attribute, handle it separately
        if isinstance(file, bytes):
            file_name = "unknown_file"
            content = file.decode("utf-8")
            documents[file_name] = content
        else:
            # Check if file is a Text File (.txt)
            if file.name.endswith(".txt"):
                content = file.read().decode("utf-8")
                documents[file.name] = content
            # Check if file is a DOCX file
            elif file.name.endswith(".docx"):
                doc = DocxDocument(io.BytesIO(file.read()))  # Use io.BytesIO to read the file content
                full_text = [para.text for para in doc.paragraphs]
                documents[file.name] = '\n'.join(full_text)
            # Check if file is a PDF file
            elif file.name.endswith(".pdf"):
                with pdfplumber.open(io.BytesIO(file.read())) as pdf:
                    full_text = [page.extract_text() for page in pdf.pages if page.extract_text()]
                    documents[file.name] = '\n'.join(full_text)
            # Check if file is a PowerPoint file
            elif file.name.endswith(".pptx"):
                prs = Presentation(io.BytesIO(file.read()))  # Use io.BytesIO to read the file content
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                documents[file.name] = text

    # Return the combined content as a single string (or the content of a specific file)
    # Example: You can combine the content of all files, or just return the first one.
    return list(documents.values())[0] if documents else ""
import re
def display_error_tabs(code_file_content, client, org_standards_documents, model_type, lang):
    if isinstance(code_file_content, str):
        code_file_content=code_file_content.splitlines()
    # Handle code content whether it's a list or string (assuming it's a list of lines)
    code_file_content = "\n".join(f"{i + 1}: {line}" for i, line in enumerate(code_file_content))

    # Track errors and suggestions
    error_counts = {}
    error_details = {}
    suggestions = ""

    # Calculate new errors and suggestions
    error_counts, error_details = calculate_errors(
        code_file_content, client, model_type, lang
    )
    suggestions = generate_suggestions(code_file_content, client, model_type)

    # Define the relevant errors based on the language
    relevant_error_types = get_relevant_error_types(lang)

    # Create a structured dictionary to return errors and suggestions
    tabs = []
    for error_type in relevant_error_types:
        error_count = error_counts.get(f"{error_type} Errors", 0)
        error_content = error_details.get(f"{error_type} Errors", "No errors found.")
        tabs.append({
            "title": f"{error_type} Errors",
            "content": error_content,
            "count": error_count
        })

    # Always add a tab for the suggested code (no count needed)
    tabs.append({
        "title": "Suggested Code",
        "content": suggestions,
        "count": None
    })

    # Return the errors and suggestions as structured data for the frontend
    return {
        "errors": error_counts,  # This will be used to calculate severity
        "error_tabs": tabs  # This will be used to display errors in the React ErrorTabs component
    }


def calculate_errors(code_file_content, client, model_type, lang):
    error_mapping = {
        "py": ["Syntax", "Runtime", "Logical", "Validation"],
        "js": ["Syntax", "Runtime", "Logical", "Validation"],
        "java": ["Syntax", "Runtime", "Logical", "Compilation"],
        "cpp": ["Syntax", "Runtime", "Logical", "Compilation"],
        "html": ["Syntax", "Validation"],
        "css": ["Syntax", "Validation"]
    }

    # Get relevant error types for the detected language
    relevant_errors = error_mapping.get(lang, [])

    error_counts = {}
    error_details = {}

    # Collect errors for relevant error types
    for error_type in relevant_errors:
        prompt = (
            f"You are a highly experienced programming expert tasked with identifying and fixing all {error_type.lower()} errors "
            "in the provided code. Please carefully review the code and, for each error found, provide the following details in a structured format. "
            "If no errors are found, simply respond with 'No errors.'\n\n"
            "For each identified error, please follow this format in separate lines and separate each error clearly:\n\n"
            "1. **Line No:** [Provide the exact line number where the error occurs]\n"
            "2.  **Error Line:** [Include the full line of code where the error is found]\n"
            "3. **Description:** [Provide a clear explanation of the error and why it is a problem]\n"
            "4. **Suggestion:** [Provide a specific and concise corrected version of the error line. Ensure this suggestion can replace the erroneous line]\n"
            "5. **Explanation:** [Briefly explain why your suggestion solves the problem and improves the code]\n\n"
            f"Here is the code for your review:\n\n{code_file_content}\n"
            "---------------------\n"
            "Make sure your response is well-structured and clear."
        )

        # Call the client to get error details from the model
        response = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model=model_type, **sampling_params
        )

        # Store the error response for this error type
        errors = response.choices[0].message.content

        error_details[f"{error_type} Errors"] = errors

        # Count the number of errors using 'Line No: ' as a marker
        error_count = len(re.findall(r"Line No:", errors))
        error_counts[f"{error_type} Errors"] = error_count

    return error_counts, error_details


def generate_suggestions(code_file_content, client, model_type):
    if code_file_content:
        suggested_prompt = (
            "Review the following code for any errors and areas of improvement. "
            "For each identified error, provide a clear explanation and suggest a corrected version of the code. "
            "Your response should be structured and easy to follow.\n\n"
            f"Code:\n{code_file_content}"
        )

        # Call the client to generate suggestions from the model
        suggested_response = client.chat.completions.create(
            messages=[{"role": "user", "content": suggested_prompt}],
            model=model_type, **sampling_params
        )
        return suggested_response.choices[0].message.content

    return "No suggestions available."

def get_relevant_error_types(lang):
    """Helper function to return relevant error types for the given language."""
    error_mapping = {
        "py": ["Syntax", "Runtime", "Logical", "Validation"],
        "js": ["Syntax", "Runtime", "Logical", "Validation"],
        "java": ["Syntax", "Runtime", "Logical", "Compilation"],
        "cpp": ["Syntax", "Runtime", "Logical", "Compilation"],
        "html": ["Syntax", "Validation"],
        "css": ["Syntax", "Validation"]
    }

    return error_mapping.get(lang, [])
from Bot.prompt import review,review_with_old
def handle_reviews(file_content,org_standards_content,client,model_type,lang,display_path):
    # Full review (complete review)
    full_review = review("complete", file_content, org_standards_content, client, model_type)
    score_data = calculate_score(org_standards_content, file_content, client, model_type)
    score = score_data["score"]
    score_explanation = score_data["explanation"]
    error_tabs_data = display_error_tabs(file_content, client, org_standards_content, model_type, lang)

    # Severity calculation for full review
    total_score = calculate_severity(error_tabs_data['errors'])
    severity_result = severity(error_tabs_data['errors'], file_content)
    summary_review = review("summary", file_content, org_standards_content, client, model_type)
    # Structure the full review response
    full_review_data = {
        "display_path": display_path,
        "review_output": full_review,
        "sumreview_output" : summary_review,
        "error_output": error_tabs_data['error_tabs'],  # Errors and suggestions in tab format
        "score": {
            "value": score,
            "explanation": score_explanation
        },
        "severity": severity_result
    }

    return full_review_data

def fetch_file_content(repo, filename, ref):
    try:
        file_content = repo.get_contents(filename, ref=ref)
        return file_content.decoded_content.decode('utf-8')
    except:
        return ""

def handle_file_status(old_content_str, new_content_str):
    flag = False  # Initialize flag at the start of the function
    diff_json = None
    if new_content_str and old_content_str:
        file_status="File Modified"
        diff_json = generate_diff_dataframe(old_content_str.splitlines(), new_content_str.splitlines())

    elif not old_content_str:
        file_status="File Added"

    elif not new_content_str:
        file_status="File Deleted"
        flag=True
        return file_status,diff_json,flag

    return file_status,diff_json,flag


import requests
import base64
from urllib.parse import urlparse

def extract_ado_info_from_url(ado_url):
    parsed_url = urlparse(ado_url)

    # Check if URL is in Azure DevOps format
    if parsed_url.netloc == 'dev.azure.com':
        path_parts = parsed_url.path.strip('/').split('/')
        if len(path_parts) >= 4:
            org = path_parts[0]
            project = path_parts[1]
            repo_name = path_parts[3]
            return org, project, repo_name
    return None, None, None


# Process the Azure DevOps repository
# Define the allowed file extensions
ALLOWED_EXTENSIONS = {'.py', '.js', '.html', '.cpp', '.css', '.java'}
def process_ado_repo(pat, organization, project, repo_name):
    ado_base_url = f"https://dev.azure.com/{organization}/{project}/_apis/git/repositories/{repo_name}/items"
    headers = {
        'Authorization': f'Basic {base64.b64encode(f":{pat}".encode()).decode()}',
    }

    # Get the list of files from the repo (this will fetch the root directory by default)
    params = {
        "scopePath": "/",  # Adjust this to fetch files from the root folder
        "recursionLevel": "Full",  # Recursively get all files in the repository
        "includeContentMetadata": "true",
        "api-version": "6.0"  # API version for Azure DevOps
    }

    response = requests.get(ado_base_url, headers=headers, params=params)

    if response.status_code == 200:
        files_data = response.json().get('value', [])
        file_paths = [file['path'] for file in files_data if not file.get('isFolder', False) and os.path.splitext(file['path'])[1] in ALLOWED_EXTENSIONS]
        return file_paths
    else:
        raise Exception(
            f"Failed to fetch files from Azure DevOps repository. Error: {response.status_code}, {response.text}")


def get_authenticated_user_email(ado_pat):
    """
    Fetch the email address of the authenticated user using the Azure DevOps PAT.
    """
    url = "https://app.vssps.visualstudio.com/_apis/profile/profiles/me?api-version=6.0"
    headers = {
        "Authorization": f"Basic {base64.b64encode(f':{ado_pat}'.encode()).decode()}",
        "Content-Type": "application/json"
    }

    print(f"Making request to: {url}")
    print(f"Authorization Header: {headers['Authorization'][:10]}... (truncated)")

    response = requests.get(url, headers=headers)

    print(f"Response Status Code: {response.status_code}")
    print(f"Response Content: {response.text}")

    if response.status_code == 401:
        raise Exception("Unauthorized: Ensure your PAT has 'Profile (Read)' scope and is encoded correctly.")
    elif response.status_code != 200:
        raise Exception(f"Failed to fetch user profile: {response.status_code} {response.text}")

    user_data = response.json()
    email = user_data.get('emailAddress')
    if not email:
        raise Exception("Email address not found in user profile data.")

    return email



def get_reviewer_id(pat, organization, project, repo, pr_number,unique_name):
    """Fetch the reviewers for a pull request."""
    api_url = f"https://dev.azure.com/{organization}/{project}/_apis/git/repositories/{repo}/pullRequests/{pr_number}/reviewers?api-version=6.0"
    
    headers = {
        'Content-Type': 'application/json'
    }

    # Make the GET request
    response = requests.get(api_url, headers=headers, auth=HTTPBasicAuth('', pat))

    if response.status_code == 200:
        reviewers = response.json()['value']
        for reviewer in reviewers:
            if reviewer['uniqueName'].lower() == unique_name.lower():
                return reviewer['id']
        
    else:
        print(f"Failed to fetch reviewers. Status code: {response.status_code}, Response: {response.text}")
        return None
def add_pr_comment(pat, organization, project, repo, pr_number, comment):
    """Add a comment to a pull request."""
    api_url = f"https://dev.azure.com/{organization}/{project}/_apis/git/repositories/{repo}/pullRequests/{pr_number}/threads?api-version=6.0"

    headers = {
        'Content-Type': 'application/json',
        'Authorization': f'Basic {base64.b64encode(f":{pat}".encode()).decode()}'
    }

    payload = {
        "comments": [
            {
                "content": comment,
                "commentType": 1  # Regular comment
            }
        ],
        "status": "active"
    }

    response = requests.post(api_url, json=payload, headers=headers)

    if response.status_code not in [200, 201]:
        raise Exception(f"Failed to add comment to PR #{pr_number}. Status code: {response.status_code}, Response: {response.text}")

def update_pr_vote(pat, organization, project, repo, pr_number, reviewer_id, status):
    """Update the vote for a pull request."""
    api_url = f"https://dev.azure.com/{organization}/{project}/_apis/git/repositories/{repo}/pullRequests/{pr_number}/reviewers/{reviewer_id}?api-version=6.0"
    
    headers = {
        'Content-Type': 'application/json'
    }

    # Vote: 10 for approve, -10 for reject
    vote = 10 if status.lower() == 'approve' else -10 if status.lower() == 'reject' else 0

    payload = {"vote": vote}

    # Make the PUT request
    response = requests.put(api_url, json=payload, headers=headers, auth=HTTPBasicAuth('', pat))

    if response.status_code in [200, 204]:
        print(f"PR #{pr_number} has been {status}d successfully.")
    else:
        print(f"Failed to {status} PR #{pr_number}. Status code: {response.status_code}, Response: {response.text}")


def complete_pull_request(pr_number, ado_url, ado_pat):
    """
    Complete (merge) a pull request in Azure DevOps after ensuring all required reviewers have approved it.
    """
    import base64
    import requests
    
    # Ensure correct API endpoint
    api_url = ado_url.replace('_git', '_apis/git/repositories')
    url = f"{api_url}/pullRequests/{pr_number}?api-version=6.0"

    # Headers with PAT
    headers = {
        "Authorization": f"Basic {base64.b64encode(f':{ado_pat}'.encode()).decode()}",
        "Content-Type": "application/json"
    }

    # Fetch the PR details
    pr_response = requests.get(url, headers=headers)
    if pr_response.status_code != 200:
        raise Exception(f"Failed to fetch PR details: {pr_response.text}")

    pr_data = pr_response.json()

    # Check if all required reviewers have approved
    reviewers = pr_data.get("reviewers", [])
    required_approvals = [
        reviewer for reviewer in reviewers if reviewer.get("isRequired")
    ]
    not_approved = [
        reviewer["displayName"]
        for reviewer in required_approvals
        if reviewer.get("vote") != 10  # 10 indicates approval
    ]

    if not_approved:
        raise Exception(
            f"PR cannot be completed. The following required reviewers have not approved:\n" + " ;\n".join(f"{reviewer}" for reviewer in not_approved)
        )



    # Fetch the latest commit ID
    last_commit_id = pr_data.get("lastMergeSourceCommit", {}).get("commitId")
    if not last_commit_id:
        raise Exception("Last merge source commit ID not found.")

    # Prepare payload to complete the PR
    payload = {
        "status": "completed",
        "lastMergeSourceCommit": {"commitId": last_commit_id}
    }

    # Complete the PR
    response = requests.patch(url, json=payload, headers=headers)
    if response.status_code not in [200, 204]:
        raise Exception(f"Failed to complete PR #{pr_number}: {response.text}")

    return True




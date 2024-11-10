# functions.py
import difflib
from docx import Document as DocxDocument
import pdfplumber
import re
import io
import pandas as pd
from pptx import Presentation
sampling_params = {
    "temperature": 0.2,       # More deterministic responses for precision
    "top_p": 0.9,             # Allows for diverse but coherent outputs
    "n": 1,                   # Generate one review
    "frequency_penalty": 0.7, # Stronger reduction of repetition
    "presence_penalty": 0.4,  # Encourages introduction of new ideas
}
def load_documents_from_files(files):
    documents = {}
    for file in files:
        # Check if file is a Text File (.txt)
        if file.name.endswith(".txt"):
            content = file.read().decode("utf-8")
            documents[file.name] = content
        # Check if file is a DOCX file
        elif file.name.endswith(".docx"):
            doc = DocxDocument(io.BytesIO(file.read()))  # Use io.BytesIO to read the file content
            full_text = [para.text for para in doc.paragraphs]
            documents[file.name] = '\n'.join(full_text)
        # Check if file is a PDF file
        elif file.name.endswith(".pdf"):
            with pdfplumber.open(io.BytesIO(file.read())) as pdf:
                full_text = [page.extract_text() for page in pdf.pages if page.extract_text()]
                documents[file.name] = '\n'.join(full_text)
        # Check if file is a PowerPoint file
        elif file.name.endswith(".pptx"):
            prs = Presentation(io.BytesIO(file.read()))  # Use io.BytesIO to read the file content
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            documents[file.name] = text

    # Return the combined content as a single string (or the content of a specific file)
    # Example: You can combine the content of all files, or just return the first one.
    return list(documents.values())[0] if documents else ""


def extract_changed_code(old_code, new_code):
    """Extract modified lines and clean them for review."""
    old_code_lines = old_code.splitlines(keepends=True)
    new_code_lines = new_code.splitlines(keepends=True)

    # Calculate the diff
    diff = list(difflib.unified_diff(old_code_lines, new_code_lines, lineterm=''))
    return diff  # Return the raw diff output


def generate_diff_dataframe(old_content, new_content):
    # Use difflib to generate the unified diff
    diff_lines = list(difflib.unified_diff(old_content, new_content, lineterm=''))

    # Track line numbers in old and new content
    old_line_num = 1
    new_line_num = 1
    diff_data = []

    for line in diff_lines:
        if line.startswith('---') or line.startswith('+++'):
            continue  # Skip version header lines
        elif line.startswith('@@'):
            parts = line.split(' ')
            old_line_num = int(parts[1].split(',')[0].replace('-', ''))
            new_line_num = int(parts[2].split(',')[0].replace('+', ''))
        elif line.startswith('-'):
            # Line removed from old content
            diff_data.append({
                'Old Line': f'{old_line_num}: {line[1:]}',  # Strip the '-'
                'New Line': '',
                'Change': 'Removed'
            })
            old_line_num += 1
        elif line.startswith('+'):
            # Line added to new content
            diff_data.append({
                'Old Line': '',
                'New Line': f'{new_line_num}: {line[1:]}',  # Strip the '+'
                'Change': 'Added'
            })
            new_line_num += 1
        elif line.startswith(' '):
            # Unchanged line
            old_line_num += 1
            new_line_num += 1

    # Create a DataFrame from the diff data
    df_diff = pd.DataFrame(diff_data)

    # Scan for matching line numbers in Old Line and New Line columns
    for i, row in df_diff.iterrows():
        old_line = row['Old Line'].split(':')[0] if row['Old Line'] else None
        for j, inner_row in df_diff.iterrows():
            new_line = inner_row['New Line'].split(':')[0] if inner_row['New Line'] else None
            # If old line number matches new line number, combine them
            if old_line and new_line and old_line == new_line:
                # Combine the two rows
                df_diff.at[i, 'New Line'] = inner_row['New Line']
                df_diff.at[i, 'Change'] = 'Modified'
                # Drop the row that was combined
                df_diff.drop(j, inplace=True)
                break  # Exit inner loop since we found a match

    # Reset the index after dropping rows
    df_diff.reset_index(drop=True, inplace=True)

    # Style the DataFrame with colors for 'Added', 'Removed', and 'Modified'
    def highlight_diff(row):
        color = ''
        if row['Change'] == 'Added':
            color = 'background-color: lightgreen'
        elif row['Change'] == 'Removed':
            color = 'background-color: lightcoral'
        elif row['Change'] == 'Modified':
            color = 'background-color: lightblue'
        return [color] * len(row)

    # Apply the style to the DataFrame
    df_styled = df_diff.style.apply(highlight_diff, axis=1)

    return df_styled


def compare_code(old_code, new_code):
    old_code_lines = old_code.splitlines(keepends=True)
    new_code_lines = new_code.splitlines(keepends=True)
    
    diff = list(difflib.unified_diff(old_code_lines, new_code_lines, lineterm=''))
    
    return ("Modified", diff) if len(diff) > 2 else ("New", new_code_lines)

import re
import streamlit as st

def display_error_tabs(code_file_content, client, org_standards_documents, model_type, lang):
    # Handle code content whether it's a list or string
    code_file_content = "\n".join(f"{i + 1}: {line}" for i, line in enumerate(code_file_content))

    # If new code is uploaded, reset session state and calculate errors
    if 'last_code' not in st.session_state or code_file_content != st.session_state.last_code:
        st.session_state.last_code = code_file_content
        st.session_state.error_counts = {}
        st.session_state.error_details = {}
        st.session_state.suggestions = ""

        # Calculate new errors and suggestions
        st.session_state.error_counts, st.session_state.error_details = calculate_errors(
            code_file_content, client, model_type, lang
        )
        st.session_state.suggestions = generate_suggestions(code_file_content, client, model_type)

    # Define the relevant errors based on the language
    relevant_error_types = get_relevant_error_types(lang)

    # Create tabs for all relevant error types, even if the count is zero
    tabs = []
    for error_type in relevant_error_types:
        error_count = st.session_state.error_counts.get(f"{error_type} Errors", 0)
        error_content = st.session_state.error_details.get(f"{error_type} Errors", "No errors found.")
        tabs.append((f"{error_type} Errors", error_content, error_count))

    # Always add a tab for the suggested code (no count needed)
    tabs.append(("Suggested Code", st.session_state.suggestions, None))

    # Display the errors in respective tabs
    with st.expander("Errors in New Code", expanded=False):
        # Create tab labels with the error count for error tabs
        tab_labels = [f"{title} ({count})" if count is not None else title for title, _, count in tabs]
        tab_instances = st.tabs(tab_labels)

        # Display error content in each tab
        for i, (title, content, _) in enumerate(tabs):
            display_error_content(tab_instances[i], title, content)

    return st.session_state.error_counts


def calculate_errors(code_file_content, client, model_type, lang):
  
    error_mapping = {
        "py": ["Syntax", "Runtime", "Logical", "Validation"],
        "js": ["Syntax", "Runtime", "Logical", "Validation"],
        "java": ["Syntax", "Runtime", "Logical", "Compilation"],
        "cpp": ["Syntax", "Runtime", "Logical", "Compilation"],
        "html": ["Syntax", "Validation"],
        "css": ["Syntax", "Validation"]
    }

    # Get relevant error types for the detected language
    relevant_errors = error_mapping.get(lang, [])

    error_counts = {}
    error_details = {}

    # Collect errors for relevant error types
    for error_type in relevant_errors:
        prompt = (
    f"You are a highly experienced programming expert tasked with identifying and fixing all {error_type.lower()} errors "
    "in the provided code. Please carefully review the code and, for each error found, provide the following details in a structured format. "
    "If no errors are found, simply respond with 'No errors.'\n\n"
    "For each identified error, please follow this format in separate lines and separate each error clearly:\n\n"
    "1. **Line No:** [Provide the exact line number where the error occurs]\n"
    "2.  **Error Line:** [Include the full line of code where the error is found]\n"
    "3. **Description:** [Provide a clear explanation of the error and why it is a problem]\n"
    "4. **Suggestion:** [Provide a specific and concise corrected version of the error line. Ensure this suggestion can replace the erroneous line]\n"
    "5. **Explanation:** [Briefly explain why your suggestion solves the problem and improves the code]\n\n"
    f"Here is the code for your review:\n\n{code_file_content}\n"
    "---------------------\n"
    "Make sure your response is well-structured and clear."
)


        # Call the client to get error details from the model
        response = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model=model_type
        )

        # Store the error response for this error type
        errors = response.choices[0].message.content
        error_details[f"{error_type} Errors"] = errors

        # Count the number of errors using 'Line No: ' as a marker
        error_count = len(re.findall(r"Line No:", errors))
        error_counts[f"{error_type} Errors"] = error_count

    return error_counts, error_details


def generate_suggestions(code_file_content, client, model_type):

    if code_file_content:
        suggested_prompt = (
            "Review the following code for any errors and areas of improvement. "
            "For each identified error, provide a clear explanation and suggest a corrected version of the code. "
            "Your response should be structured and easy to follow.\n\n"
            f"Code:\n{code_file_content}"
        )

        # Call the client to generate suggestions from the model
        suggested_response = client.chat.completions.create(
            messages=[{"role": "user", "content": suggested_prompt}],
            model=model_type
        )
        return suggested_response.choices[0].message.content

    return "No suggestions available."


def display_error_content(tab, title, content):
    with tab:
        if content:  # Ensure there's content to display
            st.write(f"### {title}")
            st.write(content)
        else:
            st.write(f"No {title.lower()} found.")


def get_relevant_error_types(lang):
    """Helper function to return relevant error types for the given language."""
    error_mapping = {
        "py": ["Syntax", "Runtime", "Logical", "Validation"],
        "js": ["Syntax", "Runtime", "Logical", "Validation"],
        "java": ["Syntax", "Runtime", "Logical", "Compilation"],
        "cpp": ["Syntax", "Runtime", "Logical", "Compilation"],
        "html": ["Syntax", "Validation"],
        "css": ["Syntax", "Validation"]
    }
    return error_mapping.get(lang, [])


def calculate_score(org_std_text, new_code,client,model_type):
    # Define prompts based on organizational code standards
    explain_prompt = (
        "Evaluate the following code according to the organization standards. "
        "Provide an overall score out of 10 in the format 'Overall Score: (score)/10' and concisely explain the reasoning behind your score.\n\n"
        f"Code:\n{new_code}\n\n"
        f"Organization Standards:\n{org_std_text}"
    )
    explain_response = client.chat.completions.create(
                        messages=[{"role": "user", "content": explain_prompt}],
                        model=model_type,**sampling_params
                    )
                    
    explanation = explain_response.choices[0].message.content
    # Assuming the response is a score as a float, parse it
    score_match = re.search(r"Overall Score.*?(\d+(\.\d+)?)/10", explanation)  # This will match a number (integer or float)

    if score_match:
        score = float(score_match.group(1))  # Extract the score and convert it to a float
    else:
        st.error("Failed to parse the score from the response.")
        return None, None

    # Display color based on score
    if score >= 9:
        color = "green"
        message = "Excellent"
    elif score >= 7:
        color = "yellow"
        message = "Good"
    elif score >= 4:
        color = "orange"
        message = "Average"
    else:
        color = "red"
        message = "Poor"

    # Display the score with color
    score_colored = f'<span style="color:{color}; font-size: 20px;">{score}</span>'
    message_colored = f'<span style="color:{color}; font-size: 20px;">{message}</span>'
    st.subheader('Overall score based on standards')
    st.markdown(f'Overall Score: {score_colored} - {message_colored}', unsafe_allow_html=True)

    return explanation  # Ensure both are returned

error_weights = {
    "Syntax Errors": 2,
    "Runtime Errors": 2,
    "Logical Errors": 1,
    "Validation Errors": 1,
    "Compilation Errors": 2
}

# Function to calculate total severity score from the dictionary {error_type: count}
def calculate_severity(error_count):
    total_score = 0
    
    # Iterate over the error_count dictionary
    for error_type, count in error_count.items():
        # Check if the error_type exists in the predefined weights
        if error_type in error_weights:
            # Multiply the count by the weight of the error type and add it to the total score
            total_score += count * error_weights[error_type]
            
    return total_score

# Determine severity level based on the total score
def determine_severity_from_score(total_score):
    if total_score == 0:
        color = "darkgreen"
        message = "No errors, perfect!"
    elif total_score <= 10:
        color = "lightgreen"
        message =  "Low severity"
    elif 11 <= total_score <= 20:
        color = "yellow"
        message = "Medium severity"
    elif 21 <= total_score <= 30:
        color = "orange"
        message = "High severity"
    else:
        color = "red"
        message = "Critical severity"
    

    return color, message
def detect_vulnerabilities(code):
    vulnerabilities = []

    # Example checks for vulnerabilities
    if "eval(" in code or "exec(" in code:
        vulnerabilities.append("Avoid using eval() or exec() as they can lead to code injection vulnerabilities.")

    if re.search(r"(\bSELECT\b|\bUPDATE\b|\bDELETE\b|\bINSERT\b)\s+\w+\s+FROM\s+\w+\s*;", code, re.IGNORECASE):
        vulnerabilities.append("Possible SQL Injection vulnerability detected. Use parameterized queries instead.")

    # Check for common XSS patterns
    if re.search(r"<script.*?>", code, re.IGNORECASE):
        vulnerabilities.append("Potential Cross-Site Scripting (XSS) vulnerability detected. Avoid using inline scripts.")

    # Check for hardcoded credentials (this is a simplistic check)
    if re.search(r"(password|secret|api_key)\s*=\s*['\"].*['\"]", code, re.IGNORECASE):
        vulnerabilities.append("Hardcoded credentials detected. Consider using environment variables or a secure vault.")

    return vulnerabilities

def severity(error_count,total_score,code_to_review):
    total_score = calculate_severity(error_count)  # Calculate total severity score
    color, severity = determine_severity_from_score(total_score)  # Determine the severity level

    # Display the score with color
    st.markdown(f"### Severity Analysis")
    st.write(f"Total Errors: {sum(error_count.values())}")
    total_score_colored = f'<span style="color:{color}; font-size: 20px;">{total_score}</span>'
    message_colored = f'<span style="color:{color}; font-size: 20px;">{severity}</span>'
    st.markdown(f'Total Severity Score: {total_score_colored } - Severity: { message_colored }', unsafe_allow_html=True)

    vulnerabilities_found = detect_vulnerabilities(code_to_review)
    # Display vulnerabilities and improvements
    if vulnerabilities_found:
        st.subheader("Detected Vulnerabilities")
        for vulnerability in vulnerabilities_found:
            st.markdown(f"- {vulnerability}")
    else:
        st.success("No vulnerabilities detected.")  
